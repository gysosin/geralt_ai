"""
Document Extractors

Extractors for common document formats: PDF, DOCX, PPTX, Excel.
Each extractor provides consistent logging and error handling.
"""
import io
import string
import logging
import pandas as pd
from typing import List, Dict, Any, Union
from PIL import Image
import pytesseract

from core.extraction.base import BaseExtractor
from core.extraction.factory import ExtractorFactory

logger = logging.getLogger(__name__)


# -------------------------------------------------------------------------
# PDF Extractor
# -------------------------------------------------------------------------
class PDFExtractor(BaseExtractor):
    """
    Extract text and images from PDF documents.
    
    Features:
    - Text extraction with bounding box coordinates
    - OCR fallback for scanned/image-based PDFs
    - Page snapshot generation for visual reference
    """
    
    def __init__(self):
        super().__init__()
    
    def extract(self, file: Union[str, bytes, io.BytesIO], **kwargs) -> List[Dict[str, Any]]:
        import fitz  # PyMuPDF
        
        self._log_start("PDF document")

        if isinstance(file, bytes):
            stream = io.BytesIO(file)
        elif isinstance(file, io.BytesIO):
            stream = file
        else:
            stream = open(file, "rb")

        extracted_pages = []
        doc = None
        
        try:
            # Open PDF
            if isinstance(stream, io.BytesIO):
                file_bytes = stream.read()
                doc = fitz.open(stream=file_bytes, filetype="pdf")
            else:
                doc = fitz.open(stream.name)
            
            total_pages = len(doc)
            self.logger.info(f"Processing PDF with {total_pages} pages")
            
            # Zoom factor for rendering (2x for quality)
            zoom_matrix = fitz.Matrix(2, 2)
            
            for page_num, page in enumerate(doc):
                self._log_debug(f"Processing page {page_num + 1}/{total_pages}")
                
                # 1. Render page image (snapshot)
                pix = page.get_pixmap(matrix=zoom_matrix)
                img_bytes = pix.tobytes("png")
                img_width, img_height = pix.width, pix.height
                
                # 2. Extract text blocks with bbox coordinates
                text_dict = page.get_text("dict")
                blocks = text_dict.get("blocks", [])
                
                # 3. OCR fallback if no text blocks
                text_blocks = [b for b in blocks if b.get("type") == 0]
                use_ocr = len(text_blocks) == 0 or all(
                    len(b.get("lines", [])) == 0 for b in text_blocks
                )
                
                if use_ocr:
                    self.logger.info(f"Page {page_num + 1}: No text detected, attempting OCR")
                    try:
                        pil_img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                        ocr_text = pytesseract.image_to_string(pil_img)
                        if ocr_text.strip():
                            self._log_debug(f"Page {page_num + 1}: OCR extracted {len(ocr_text)} characters")
                            blocks = [{
                                "type": 0,
                                "bbox": [0, 0, page.rect.width, page.rect.height],
                                "lines": [{
                                    "spans": [{
                                        "text": ocr_text,
                                        "bbox": [0, 0, page.rect.width, page.rect.height]
                                    }]
                                }]
                            }]
                        else:
                            self._log_warning(f"Page {page_num + 1}: OCR returned no text")
                    except Exception as ocr_error:
                        self._log_warning(f"Page {page_num + 1}: OCR failed - {ocr_error}")
                
                # 4. Process text blocks into chunks
                first_block_of_page = True
                block_idx = 0
                page_blocks_count = 0
                
                for block in blocks:
                    if block.get("type") != 0:  # Skip image blocks
                        continue
                    
                    # Extract all text from this block
                    block_text = ""
                    block_bbox = block.get("bbox", [0, 0, page.rect.width, page.rect.height])
                    
                    for line in block.get("lines", []):
                        for span in line.get("spans", []):
                            block_text += span.get("text", "") + " "
                        block_text += "\n"
                    
                    block_text = block_text.strip()
                    if not block_text:
                        continue
                    
                    block_idx += 1
                    page_blocks_count += 1
                    
                    # Convert PDF coordinates to image coordinates (apply zoom)
                    scale_x = zoom_matrix.a
                    scale_y = zoom_matrix.d
                    
                    bbox_scaled = [
                        int(block_bbox[0] * scale_x),
                        int(block_bbox[1] * scale_y),
                        int(block_bbox[2] * scale_x),
                        int(block_bbox[3] * scale_y)
                    ]
                    
                    item = {
                        "content": block_text,
                        "metadata": {
                            "page_number": page_num + 1,
                            "block_index": block_idx,
                            "bbox": bbox_scaled,
                            "image_dimensions": {
                                "width": img_width,
                                "height": img_height
                            }
                        },
                    }
                    
                    # Attach snapshot to first block of page
                    if first_block_of_page:
                        item["_page_image_bytes"] = img_bytes
                        first_block_of_page = False
                    
                    extracted_pages.append(item)
                
                self._log_debug(f"Page {page_num + 1}: Extracted {page_blocks_count} text blocks")

            doc.close()
            if hasattr(stream, 'close') and stream is not file:
                stream.close()
            
            self._log_complete(len(extracted_pages), "text blocks")
            return extracted_pages

        except Exception as e:
            self._log_error(e, "PDF processing")
            if doc:
                doc.close()
            raise ValueError(f"Error processing PDF: {str(e)}")


ExtractorFactory.register("pdf", PDFExtractor)


# -------------------------------------------------------------------------
# DOCX Extractor
# -------------------------------------------------------------------------
class DocxExtractor(BaseExtractor):
    """
    Extract text, tables, and OCR'd images from DOCX documents.
    
    Features:
    - Paragraph extraction with style information
    - Table extraction with header-value pairing
    - Image OCR for embedded images
    """
    
    def __init__(self):
        super().__init__()
    
    def extract(self, file: Union[str, bytes, io.BytesIO], **kwargs) -> List[Dict[str, Any]]:
        from docx import Document

        self._log_start("DOCX document")
        
        if isinstance(file, bytes):
            file = io.BytesIO(file)

        from_page = kwargs.get("from_page", 0)
        to_page = kwargs.get("to_page", 999999999)
        do_ocr = kwargs.get("do_ocr", True)

        results = []
        page_number = 0
        paragraph_counter = 0
        table_counter = 0

        try:
            doc = Document(file)
            
            total_paragraphs = len(doc.paragraphs)
            total_tables = len(doc.tables)
            self.logger.info(f"DOCX contains {total_paragraphs} paragraphs and {total_tables} tables")

            # 1) Paragraph text
            for para in doc.paragraphs:
                if page_number >= to_page:
                    break
                    
                runs_text = []
                for run in para.runs:
                    if "lastRenderedPageBreak" in run._element.xml:
                        page_number += 1
                        if page_number >= to_page:
                            break

                    if from_page <= page_number < to_page and run.text.strip():
                        runs_text.append(run.text)

                merged = "".join(runs_text).strip()
                if merged:
                    paragraph_counter += 1
                    results.append({
                        "content": merged,
                        "metadata": {
                            "page_number": page_number,
                            "paragraph_number": paragraph_counter,
                            "style_name": para.style.name if para.style else None,
                        },
                    })

            self._log_debug(f"Extracted {paragraph_counter} paragraphs")

            # 2) Tables
            for table in doc.tables:
                table_counter += 1
                row_count = len(table.rows)
                col_count = len(table.columns)
                
                if row_count == 0 or col_count == 0:
                    self._log_debug(f"Table {table_counter}: Skipping empty table")
                    continue

                headers = []
                for c_idx in range(col_count):
                    text_h = table.cell(0, c_idx).text.strip()
                    headers.append(text_h)

                rows_extracted = 0
                for row_idx in range(1, row_count):
                    row_vals = []
                    for c_idx in range(col_count):
                        cell_val = table.cell(row_idx, c_idx).text.strip()
                        header_val = headers[c_idx] if c_idx < len(headers) else ""
                        if header_val or cell_val:
                            row_vals.append(f"{header_val}: {cell_val}".strip(": "))
                    if row_vals:
                        rows_extracted += 1
                        results.append({
                            "content": "; ".join(row_vals),
                            "metadata": {
                                "page_number": page_number,
                                "table_number": table_counter,
                                "row_number": row_idx,
                            },
                        })
                
                self._log_debug(f"Table {table_counter}: Extracted {rows_extracted} rows")

            # 3) Images + OCR
            if do_ocr:
                image_count = 0
                ocr_success = 0
                
                for rel in doc.part.rels.values():
                    if "image" in rel.target_ref:
                        image_count += 1
                        try:
                            blob = rel.target_part.blob
                            pil_img = Image.open(io.BytesIO(blob)).convert("L")
                            ocr_text = pytesseract.image_to_string(pil_img)
                            if ocr_text.strip():
                                ocr_success += 1
                                paragraph_counter += 1
                                results.append({
                                    "content": ocr_text.strip(),
                                    "metadata": {
                                        "page_number": page_number,
                                        "paragraph_number": paragraph_counter,
                                        "ocr_extracted": True,
                                    },
                                })
                        except Exception as img_error:
                            self._log_warning(f"Image OCR failed for image {image_count}: {img_error}")
                
                if image_count > 0:
                    self.logger.info(f"Processed {image_count} images, OCR succeeded for {ocr_success}")

            self._log_complete(len(results), "content blocks")
            return results
            
        except Exception as e:
            self._log_error(e, "DOCX processing")
            raise ValueError(f"Error processing DOCX: {str(e)}")


ExtractorFactory.register("docx", DocxExtractor)


# -------------------------------------------------------------------------
# PPTX Extractor
# -------------------------------------------------------------------------
class PPTXExtractor(BaseExtractor):
    """
    Extract text, tables, and OCR'd images from PPTX presentations.
    
    Features:
    - Text extraction from all shape types
    - Table extraction with header-value pairing
    - Image OCR for embedded images
    - Group shape support
    """
    
    def __init__(self):
        super().__init__()
    
    def extract(self, file: Union[str, bytes, io.BytesIO], **kwargs) -> List[Dict[str, Any]]:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        self._log_start("PPTX presentation")
        
        if isinstance(file, bytes):
            file = io.BytesIO(file)

        from_slide = kwargs.get("from_slide", 0)
        to_slide = kwargs.get("to_slide", None)

        try:
            pres = Presentation(file)
            total_slides = len(pres.slides)
            to_slide = to_slide or total_slides
            
            self.logger.info(f"Processing PPTX with {total_slides} slides (range: {from_slide + 1}-{to_slide})")

            flattened = []
            
            for slide_num, slide in enumerate(pres.slides, start=1):
                if slide_num < from_slide + 1:
                    continue
                if slide_num > to_slide:
                    break

                slide_items = 0
                sorted_shapes = sorted(
                    slide.shapes, key=lambda s: ((s.top or 0) // 10, (s.left or 0))
                )
                
                for sh_idx, shape in enumerate(sorted_shapes):
                    if shape.has_text_frame:
                        txt = "\n".join(
                            p.text.strip()
                            for p in shape.text_frame.paragraphs
                            if p.text.strip()
                        )
                        if txt:
                            slide_items += 1
                            flattened.append({
                                "content": txt,
                                "metadata": {
                                    "slide_number": slide_num,
                                    "shape_type": str(shape.shape_type),
                                    "shape_index": sh_idx + 1,
                                },
                            })
                            
                    elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                        table = shape.table
                        row_count = len(table.rows)
                        col_count = len(table.columns)
                        if row_count == 0 or col_count == 0:
                            continue

                        headers = []
                        for c_idx in range(col_count):
                            headers.append(table.cell(0, c_idx).text.strip())

                        table_texts = []
                        for row_i in range(1, row_count):
                            row_vals = []
                            for c_i in range(col_count):
                                hdr = headers[c_i] if c_i < len(headers) else ""
                                cell_txt = table.cell(row_i, c_i).text.strip()
                                if hdr or cell_txt:
                                    row_vals.append(f"{hdr}: {cell_txt}".strip(": "))
                            if row_vals:
                                table_texts.append("; ".join(row_vals))
                                
                        if table_texts:
                            slide_items += 1
                            joined_tbl = "\n".join(table_texts)
                            flattened.append({
                                "content": joined_tbl,
                                "metadata": {
                                    "slide_number": slide_num,
                                    "shape_type": "table",
                                    "shape_index": sh_idx + 1,
                                },
                            })
                            
                    elif hasattr(shape, "image") and shape.image:
                        try:
                            img_bytes = io.BytesIO(shape.image.blob)
                            pil_img = Image.open(img_bytes).convert("L")
                            ocr_text = pytesseract.image_to_string(pil_img)
                            if ocr_text.strip():
                                slide_items += 1
                                flattened.append({
                                    "content": ocr_text.strip(),
                                    "metadata": {
                                        "slide_number": slide_num,
                                        "shape_type": "ocr",
                                        "shape_index": sh_idx + 1,
                                    },
                                })
                        except Exception as img_err:
                            self._log_warning(f"Slide {slide_num}: Image OCR failed - {img_err}")
                            
                    elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                        sub_texts = []
                        subshapes_sorted = sorted(
                            shape.shapes, key=lambda x: ((x.top or 0) // 10, (x.left or 0))
                        )
                        for sb in subshapes_sorted:
                            try:
                                if sb.has_text_frame:
                                    stxt = "\n".join(
                                        pp.text.strip()
                                        for pp in sb.text_frame.paragraphs
                                        if pp.text.strip()
                                    )
                                    if stxt:
                                        sub_texts.append(stxt)
                            except Exception as grp_err:
                                self._log_debug(f"Slide {slide_num}: Group shape text extraction failed - {grp_err}")
                                
                        if sub_texts:
                            slide_items += 1
                            flattened.append({
                                "content": "\n".join(sub_texts),
                                "metadata": {
                                    "slide_number": slide_num,
                                    "shape_type": "group",
                                    "shape_index": sh_idx + 1,
                                },
                            })
                
                self._log_debug(f"Slide {slide_num}: Extracted {slide_items} items")
            
            self._log_complete(len(flattened), "content items")
            return flattened
            
        except Exception as e:
            self._log_error(e, "PPTX processing")
            raise ValueError(f"Error processing PPTX: {str(e)}")


ExtractorFactory.register("pptx", PPTXExtractor)


# -------------------------------------------------------------------------
# Excel Extractor
# -------------------------------------------------------------------------
class ExcelExtractor(BaseExtractor):
    """
    Extract cell data from Excel spreadsheets.
    
    Features:
    - Multi-sheet support
    - Cell address tracking
    - Handles both .xlsx and .xls formats
    """
    
    def __init__(self):
        super().__init__()
    
    def extract(self, file: Union[str, bytes, io.BytesIO], **kwargs) -> List[Dict[str, Any]]:
        self._log_start("Excel spreadsheet")
        
        extracted_cells = []
        try:
            sheets = pd.read_excel(file, sheet_name=None)
            
            self.logger.info(f"Processing {len(sheets)} sheets")
            
            for sheet_name, df in sheets.items():
                sheet_cells = 0
                col_count = len(df.columns)
                
                for row_idx, row in df.iterrows():
                    for col_idx, value in row.items():
                        if pd.notna(value):
                            # Handle columns beyond Z (AA, AB, etc.)
                            col_pos = df.columns.get_loc(col_idx)
                            if col_pos < 26:
                                col_letter = string.ascii_uppercase[col_pos]
                            else:
                                # Handle columns beyond Z (AA, AB, etc.)
                                first_letter = string.ascii_uppercase[(col_pos // 26) - 1]
                                second_letter = string.ascii_uppercase[col_pos % 26]
                                col_letter = first_letter + second_letter
                            
                            cell_addr = f"{col_letter}{row_idx + 2}"
                            sheet_cells += 1
                            extracted_cells.append({
                                "content": str(value),
                                "metadata": {
                                    "cell_address": cell_addr,
                                    "sheet_name": sheet_name,
                                },
                            })
                
                self._log_debug(f"Sheet '{sheet_name}': Extracted {sheet_cells} cells")
            
            self._log_complete(len(extracted_cells), "cells")
            return extracted_cells
            
        except Exception as e:
            self._log_error(e, "Excel processing")
            raise ValueError(f"Error extracting from Excel: {str(e)}")


ExtractorFactory.register("xlsx", ExcelExtractor)
ExtractorFactory.register("xls", ExcelExtractor)
